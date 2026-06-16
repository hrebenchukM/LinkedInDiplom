import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt

ROOT_DIR = Path(__file__).resolve().parent.parent
DOCS_DIR = Path(__file__).resolve().parent
OUTPUT_DOCS_FILE = DOCS_DIR / "LinkedIn_Diploma_Presentation_UA.pptx"
OUTPUT_ROOT_FILE = ROOT_DIR / "presentation.pptx"
FS_BACKEND_IMG = DOCS_DIR / "assets" / "filesystem-backend.png"
FS_FRONTEND_IMG = DOCS_DIR / "assets" / "filesystem-frontend.png"

WHITE = RGBColor(255, 255, 255)
BG_DARK = RGBColor(7, 11, 20)
BG_DARK_2 = RGBColor(11, 18, 32)
CARD = RGBColor(18, 26, 43)
CARD_SOFT = RGBColor(27, 38, 59)
CARD_LIGHT = RGBColor(36, 50, 75)
BLUE = RGBColor(10, 102, 194)
BLUE_SOFT = RGBColor(42, 64, 96)
CYAN = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
ORANGE = RGBColor(251, 146, 60)
PURPLE = RGBColor(139, 92, 246)
PURPLE_SOFT = RGBColor(196, 181, 253)
TEXT = RGBColor(238, 242, 247)
TEXT_MUTED = RGBColor(148, 163, 184)
BORDER = RGBColor(51, 65, 85)

NAV_ITEMS = [
    ("home", "Головна"),
    ("network", "Мережа"),
    ("vacancies", "Вакансії"),
    ("chat", "Чат"),
    ("profile", "Профіль"),
]

SLIDES = [
    {
        "type": "cover",
        "focus": None,
        "title": "LinkedIn Diplom Project",
        "subtitle": "Платформа професійного нетворкінгу, пошуку роботи та комунікації",
        "facts": [
            "Дипломний проєкт: LinkedIn-like MVP",
            "Frontend: Ямчук Тимур, Андрій Ротарь",
            "Backend: Гребенчук Марія, Прохорова Валерія",
            "React/Vite + ASP.NET Core .NET 8 + PostgreSQL",
            "Рік захисту: 2026",
        ],
    },
    {
        "type": "statement",
        "focus": None,
        "title": "1. Актуальність теми",
        "headline": "Ринок потребує не просто дошки вакансій, а цілісної професійної екосистеми.",
        "points": [
            "Кандидатам потрібні профіль, мережа контактів, контент, вакансії та комунікація в одному місці.",
            "Роботодавцям важливо бачити не лише резюме, а й професійну активність користувача.",
            "Локальні job-board часто відокремлюють пошук роботи від нетворкінгу та обміну досвідом.",
        ],
        "metric": "5",
        "metric_label": "ключових сценаріїв об'єднано в одному MVP",
    },
    {
        "type": "statement",
        "focus": None,
        "title": "Обґрунтування вибору LinkedIn-like формату",
        "headline": "LinkedIn обрано як модель, бо він поєднує роботу, професійну репутацію, контент і контакти в одному продукті.",
        "points": [
            "Такий формат дозволяє показати не одну сторінку вакансій, а повноцінну систему з різними ролями користувачів.",
            "Для дипломного проєкту це сильніше за простий job-board: є frontend, backend, база даних, авторизація та інтеграція модулів.",
            "Користувацький шлях легко демонструвати на захисті: вхід, профіль, стрічка, мережа, вакансії, чат.",
        ],
        "metric": "6",
        "metric_label": "робочих розділів сайту демонструють цілісний продукт",
    },
    {
        "type": "comparison",
        "focus": None,
        "title": "Перевага LinkedIn Clone над конкурентами",
        "left_title": "Конкуренти та типові аналоги",
        "right_title": "LinkedIn Clone у цьому проєкті",
        "left": [
            "Job-board часто обмежується списком вакансій і формою відгуку.",
            "Профіль кандидата існує окремо від контенту та комунікації.",
            "Чати, спільноти й події зазвичай винесені в інші сервіси.",
            "Демо часто показує UI без повної прив'язки до backend API.",
        ],
        "right": [
            "Один сайт об'єднує вакансії, профіль, стрічку, мережу, чати та сповіщення.",
            "Профіль працює як професійна візитка з навичками, досвідом і резюме.",
            "Спільноти, сторінки, events і messaging створюють довгострокову взаємодію.",
            "Frontend інтегрований з .NET API, JWT, PostgreSQL і Swagger-контрактами.",
        ],
    },
    {
        "type": "modules",
        "focus": None,
        "title": "Ключові переваги саме нашого сайту",
        "subtitle": "Переваги сформовані з реалізованого коду та доступних екранів продукту",
        "items": [
            ("Єдиний UX", "всі сценарії в одному інтерфейсі"),
            ("Реальний API", "auth, profile, jobs, feed, messaging"),
            ("Гнучкий demo", "backend або локальні fallback-дані"),
            ("Сучасний UI", "dark/light theme, i18n, animations"),
            ("Масштабованість", "модульний моноліт з BFF-шаром"),
            ("Безпека", "JWT, refresh, ролі User/Admin"),
            ("Відтворюваність", "Docker Compose і Swagger"),
            ("Roadmap", "realtime, e2e tests, analytics"),
        ],
    },
    {
        "type": "matrix",
        "focus": None,
        "title": "Команда проєкту та зони відповідальності",
        "left_title": "Frontend розробники",
        "right_title": "Backend розробники",
        "left": [
            "Ямчук Тимур: реєстрація, вакансії, пошук роботи, чати, дизайн і особистий кабінет.",
            "Андрій Ротарь: головний екран, контент, вкладка спільноти та інтеграція frontend з backend.",
            "Frontend-команда відповідає за користувацький шлях і демонстрацію сайту.",
        ],
        "right": [
            "Гребенчук Марія: backend-частина доменних модулів, API та робота з даними.",
            "Прохорова Валерія: backend-логіка, інтеграція модулів, стабільність API та БД.",
            "Backend-команда забезпечує REST API, PostgreSQL, авторизацію та бізнес-логіку.",
        ],
    },
    {
        "type": "matrix",
        "focus": None,
        "title": "2. Мета і завдання проєкту",
        "left_title": "Мета",
        "right_title": "Завдання",
        "left": [
            "Розробити вебплатформу для професійної взаємодії.",
            "Забезпечити наскрізний шлях користувача від реєстрації до відгуку на вакансію.",
            "Показати готовність архітектури до подальшого масштабування.",
        ],
        "right": [
            "Створити SPA frontend з основними сторінками продукту.",
            "Реалізувати REST API, JWT-auth, модулі профілю, контенту, мережі, вакансій і чатів.",
            "Підготувати Docker Compose, Swagger і базові тести для перевірки рішення.",
        ],
    },
    {
        "type": "modules",
        "focus": None,
        "title": "3. Огляд функціональності продукту",
        "subtitle": "MVP покриває основні сценарії LinkedIn-подібної платформи",
        "items": [
            ("Auth", "реєстрація, login, JWT, refresh"),
            ("Profile", "дані користувача, навички, резюме"),
            ("Feed", "пости, реакції, коментарі, share"),
            ("Network", "контакти, підписки, групи, сторінки"),
            ("Jobs", "пошук, фільтри, apply, saved"),
            ("Messaging", "чати, unread, AI assistant"),
            ("Notifications", "події та прочитані стани"),
            ("Admin", "ролі, модерація, статистика"),
        ],
    },
    {
        "type": "comparison",
        "focus": None,
        "title": "4. Порівняння з ринковими аналогами",
        "left_title": "Типовий job-board",
        "right_title": "Наш LinkedIn-like MVP",
        "left": [
            "Вакансії як головний і майже єдиний сценарій.",
            "Профіль кандидата обмежений резюме.",
            "Комунікація часто переноситься в зовнішні месенджери.",
            "Мало контенту, спільнот і професійної взаємодії.",
        ],
        "right": [
            "Вакансії + профіль + мережа + контент + чат.",
            "Профіль містить навички, досвід, освіту, резюме та avatar/upload.",
            "Комунікація проходить через `/api/messaging/*` всередині платформи.",
            "Є groups, pages, events і глобальний пошук по продукту.",
        ],
    },
    {
        "type": "stack",
        "focus": None,
        "title": "5. Технологічний стек",
        "groups": [
            ("Frontend", ["React 19", "Vite 6", "React Router 7", "Context API", "CSS variables"]),
            ("Backend", ["ASP.NET Core .NET 8", "Modular Monolith", "BFF Facade.API", "Swagger/OpenAPI"]),
            ("Data & Security", ["PostgreSQL 16", "EF Core 8", "JWT Bearer", "Refresh tokens"]),
            ("Infra & QA", ["Docker Compose", "Uploads volume", "xUnit + Moq", "EF InMemory tests"]),
        ],
    },
    {
        "type": "architecture",
        "focus": None,
        "title": "6. Загальна архітектура системи",
        "steps": [
            ("React/Vite SPA", "сторінки, routing, state"),
            ("HTTP API Client", "tokens, refresh, errors"),
            ("Facade.API", "BFF controllers + services"),
            ("Core modules", "business logic + clients"),
            ("PostgreSQL", "schema per module"),
        ],
        "note": "Архітектура побудована як microservice-ready modular monolith: один процес зараз, чіткі межі модулів для майбутнього винесення.",
    },
    {
        "type": "modules",
        "focus": None,
        "title": "7. Backend: модульний моноліт + BFF",
        "subtitle": "9 core-модулів і 10 facade-модулів формують близько 200 REST endpoints",
        "items": [
            ("Identity", "користувачі, ролі, JWT"),
            ("Profile", "профілі та перегляди"),
            ("Professional", "досвід, освіта, навички"),
            ("Network", "контакти, групи, сторінки"),
            ("Content", "пости, реакції, hashtags"),
            ("Messaging", "чати та повідомлення"),
            ("Jobs", "вакансії та відгуки"),
            ("Events/Admin", "події, модерація, stats"),
        ],
    },
    {
        "type": "database",
        "focus": None,
        "title": "8. База даних і міграції",
        "schemas": ["identity", "profile", "professional", "network", "content", "messaging", "jobs", "notifications", "events"],
        "points": [
            "PostgreSQL 16 з окремою схемою для кожного доменного модуля.",
            "EF Core migrations запускаються при старті API в контрольованому порядку.",
            "Для завантажених файлів використовується `/uploads` і Docker volume.",
            "Platform Admin реалізовано через JWT role `Admin`, без окремої таблиці адміністраторів.",
        ],
    },
    {
        "type": "security",
        "focus": "profile",
        "title": "9. Безпека та авторизація",
        "points": [
            ("JWT access token", "короткоживучий токен для захищених endpoint-ів"),
            ("Refresh token", "зберігається в PostgreSQL і ротується при refresh"),
            ("Roles", "`User` для звичайних користувачів, `Admin` для `/api/admin/*`"),
            ("Claims", "userId береться з JWT, а не з даних, надісланих клієнтом"),
        ],
    },
    {
        "type": "frontend",
        "focus": "home",
        "title": "10. Frontend: сторінки та UX",
        "routes": [
            ("/auth", "реєстрація, login, social demo"),
            ("/home", "стрічка, composer, mini inbox"),
            ("/network", "контакти, групи, pages/events"),
            ("/vacancies", "вакансії, фільтри, apply"),
            ("/chat", "повідомлення, AI assistant, calls"),
            ("/profile", "профіль, avatar, skills, resume"),
        ],
    },
    {
        "type": "journey",
        "focus": "home",
        "title": "11. Наскрізний сценарій користувача",
        "steps": [
            "Реєстрація або вхід",
            "Заповнення профілю",
            "Перегляд стрічки",
            "Розширення мережі",
            "Пошук вакансій",
            "Чат з контактом",
        ],
        "note": "Цей шлях демонструє, що frontend і backend працюють як одна система, а не як окремі макети.",
    },
    {
        "type": "feature",
        "focus": "vacancies",
        "title": "12. Модуль вакансій",
        "headline": "Сценарій кандидата: знайти, зберегти, відгукнутися та відстежити заявку.",
        "points": [
            "Фільтри за ключовими словами, локацією, типом зайнятості, рівнем досвіду та remote.",
            "Підтримка saved jobs, apply flow, my applications і публікації вакансій.",
            "API-шар використовує `/api/jobs/vacancies` та пов'язані jobs endpoints.",
            "UI готовий для демо як з backend API, так і з локальними fallback-даними.",
        ],
    },
    {
        "type": "feature",
        "focus": "chat",
        "title": "13. Чати та AI-помічник",
        "headline": "Комунікація залишається всередині платформи та доповнюється локальним assistant-сценарієм.",
        "points": [
            "Чати підтримують unread states, archive/mute, видалення повідомлень і share posts.",
            "AI assistant має quick prompts, command chips і навігаційні дії.",
            "Для демо реалізовано call overlay з mute/speaker/end call станами.",
            "Backend-контракт винесено в messaging API, realtime залишено для roadmap.",
        ],
    },
    {
        "type": "feature",
        "focus": "profile",
        "title": "14. Профіль і професійні дані",
        "headline": "Профіль є центром професійної ідентичності користувача.",
        "points": [
            "Редагування профілю, avatar upload/presets, resume upload і visibility settings.",
            "Професійний блок: skills, experience, education, companies, certificates.",
            "ProfileProvider синхронізує frontend state з backend API або локальним fallback.",
            "Профіль використовується у вакансіях, мережі контактів і загальному UX.",
        ],
    },
    {
        "type": "integration",
        "focus": None,
        "title": "15. API-інтеграція frontend і backend",
        "rows": [
            ("Auth", "`/api/auth/register`, `/api/auth/login`, `/api/auth/refresh`, `/api/auth/me`"),
            ("Profile", "`/api/profile/me`, uploads, settings"),
            ("Content", "`/api/content/me/posts`, comments, reactions, saved"),
            ("Network", "`/api/network/me/contacts`, follows, groups, pages"),
            ("Jobs", "`/api/jobs/vacancies`, applications, favorites"),
            ("Messaging", "`/api/messaging/me/chats`, messages, read status"),
        ],
    },
    {
        "type": "quality",
        "focus": None,
        "title": "16. Тестування та контроль якості",
        "items": [
            ("Build", "`dotnet build LinkedIn.sln`"),
            ("Unit tests", "26 tests: ProfileService, PostService, HashtagService"),
            ("Swagger", "ручна перевірка protected endpoints з Bearer token"),
            ("QA flow", "register/login → profile → feed → network → jobs → chat"),
            ("Frontend build", "`npm run build` для Vite SPA"),
            ("Known gaps", "немає realtime та integration tests у v1"),
        ],
    },
    {
        "type": "infrastructure",
        "focus": None,
        "title": "17. Інфраструктура та запуск",
        "points": [
            ("Frontend", "http://localhost:5173, Vite dev server"),
            ("API", "http://localhost:5000/swagger у Docker або локальний `Facade.API`"),
            ("Database", "PostgreSQL 16, база `linkedin_dev`, порт 5432"),
            ("Docker Compose", "Postgres + API + frontend, healthcheck, env vars, uploads volume"),
        ],
    },
    {
        "type": "roadmap",
        "focus": None,
        "title": "18. Обмеження v1 та план розвитку",
        "steps": [
            ("Realtime", "SignalR для чатів і сповіщень"),
            ("Email flows", "verification, password reset"),
            ("QA", "integration/API tests і e2e сценарії"),
            ("Scale", "outbox/broker для domain events"),
            ("Analytics", "моніторинг, метрики, рекомендації"),
        ],
    },
    {
        "type": "roles",
        "focus": None,
        "title": "Командний внесок frontend",
        "left_title": "Ямчук Тимур",
        "right_title": "Андрій Ротарь",
        "left": [
            "Вікно реєстрації: форма входу, валідація, обробка помилок.",
            "Сторінка вакансій: список, картки, фільтри та відгук.",
            "Пошук роботи, чати, дизайн і особистий кабінет користувача.",
        ],
        "right": [
            "Головний екран: структура першого екрана та навігація.",
            "Контент на сайті: стрічка, взаємодії, візуальна подача.",
            "Вкладка спільноти та злиття backend з frontend.",
        ],
    },
    {
        "type": "feature",
        "focus": "vacancies",
        "title": "Виступ: Ямчук Тимур",
        "headline": "Ямчук Тимур презентує кандидатський шлях: від входу в систему до пошуку роботи та комунікації.",
        "points": [
            "Вікно реєстрації: структура форми, UX авторизації, валідація та повідомлення про помилки.",
            "Сторінка вакансій: картки вакансій, фільтри, збереження та підготовка сценарію відгуку.",
            "Пошук роботи: підбір за параметрами, сценарій без результатів і повторний пошук.",
            "Чати, дизайн і особистий кабінет: єдиний стиль, профіль користувача та комунікація.",
        ],
    },
    {
        "type": "feature",
        "focus": "home",
        "title": "Виступ: Андрій Ротарь",
        "headline": "Андрій Ротарь презентує соціальну частину платформи та інтеграцію frontend з backend.",
        "points": [
            "Головний екран: перший контакт користувача із сайтом, навігація та швидкий доступ до модулів.",
            "Контент на сайті: стрічка постів, взаємодії, лайки, коментарі та оновлення даних.",
            "Вкладка спільноти: контакти, groups, pages, events і розвиток професійної мережі.",
            "Злиття backend з frontend: API-клієнт, proxy, token flow і робота з реальними endpoint-ами.",
        ],
    },
    {
        "type": "roles",
        "focus": None,
        "title": "Внесок backend розробників",
        "left_title": "Гребенчук Марія",
        "right_title": "Прохорова Валерія",
        "left": [
            "Участь у backend-модулях і REST API для основних сценаріїв платформи.",
            "Робота з доменною логікою, DTO/контрактами та структурою даних.",
            "Підтримка стабільності backend-частини для демонстрації на захисті.",
        ],
        "right": [
            "Участь у backend-логіці, інтеграції модулів і роботі з PostgreSQL.",
            "Підтримка авторизації, endpoints і взаємодії frontend з API.",
            "Перевірка працездатності backend-сценаріїв через Swagger і Docker-запуск.",
        ],
    },
    {
        "type": "final",
        "focus": None,
        "title": "20. Підсумок",
        "facts": [
            "Реалізовано працездатний LinkedIn-like MVP з frontend, backend, БД та Docker-запуском.",
            "Система має модульну архітектуру, JWT-безпеку, REST API і реальні користувацькі сценарії.",
            "Проєкт готовий до демонстрації на захисті та має зрозумілий roadmap після v1.",
        ],
    },
]


def normalize_slide_numbers():
    number = 1
    for slide in SLIDES:
        if slide.get("type") == "cover":
            continue
        title = re.sub(r"^\d+\.\s*", "", slide["title"])
        slide["title"] = f"{number}. {title}"
        number += 1


PROJECT_FILESYSTEM = {
    "root": "LinkedInDiplom-master",
    "sections": [
        {
            "folder": "frontend",
            "color_idx": 2,
            "items": [
                "src/pages",
                "src/features",
                "src/app",
                "vite.config.js",
                "package.json",
            ],
            "note": "React 19 + Vite SPA — наш фокус",
        },
        {
            "folder": "backend",
            "color_idx": 4,
            "items": [
                "Identity",
                "Profile / Professional",
                "Network / Content",
                "Messaging / Jobs",
                "Events / Notifications",
                "Facade.API",
                "Tests",
            ],
            "note": "модульний моноліт — для контексту",
        },
        {
            "folder": "docs",
            "color_idx": 1,
            "items": ["architecture", "api", "E2E checklist"],
            "note": "документація проєкту",
        },
    ],
}

FRONTEND_TREE_LINES = [
    "frontend/",
    "├── src/",
    "│   ├── app/           router, providers, layout",
    "│   ├── pages/         auth, home, network, jobs, chat, admin",
    "│   ├── features/      auth, profile, jobs, chat, network",
    "│   └── shared/",
    "│       ├── api/       HTTP-клієнт, paths",
    "│       ├── ui/        спільні компоненти",
    "│       └── lib/       storage, session",
    "├── vite.config.js",
    "└── package.json",
]

BACKEND_TREE_LINES = [
    "backend/  (контекст масштабу)",
    "├── Identity · Profile · Professional",
    "├── Network · Content · Messaging",
    "├── Jobs · Events · Notifications",
    "└── Facade.API",
]


SLIDES = [
    {
        "type": "cover",
        "focus": None,
        "title": "LinkedIn Clone: LinkedIn-like платформа",
        "subtitle": "Frontend-частина дипломного проєкту професійного нетворкінгу",
        "facts": [
            "Проєкт: сайт для професійної взаємодії, пошуку роботи та комунікації",
            "Frontend розробники: Ямчук Тимур, Андрій Ротарь",
            "Стек frontend: React 19, Vite 6, React Router, Context API",
            "Demo-режим: social demo та mock-дані для стабільної демонстрації на захисті",
        ],
    },
    {
        "type": "statement",
        "focus": None,
        "title": "Вступ: чому саме LinkedIn",
        "headline": "LinkedIn обрано як найсильнішу модель, бо він поєднує професійний профіль, мережу контактів, контент, вакансії та комунікацію.",
        "points": [
            "Це не просто сайт з вакансіями, а повна професійна екосистема для кандидата, роботодавця і спільнот.",
            "Такий формат дозволяє показати більше дипломної роботи: UI, маршрути, клієнтські стани та користувацькі сценарії.",
            "Клон LinkedIn добре демонструється на захисті через demo-flow: вхід → профіль → стрічка → вакансії → чат.",
        ],
        "metric": "6",
        "metric_label": "основних розділів об'єднані в одному сайті",
    },
    {
        "type": "filesystem",
        "focus": None,
        "title": "Файлова структура проєкту",
        "subtitle": "Фокус доповіді — frontend SPA; backend показано лише для розуміння масштабу системи",
        "tree": PROJECT_FILESYSTEM,
    },
    {
        "type": "comparison",
        "focus": None,
        "title": "Порівняння з конкретними конкурентами",
        "left_title": "Що є у Work.ua, LinkedIn, Djinni, Robota.ua",
        "right_title": "Чого бракує там, але є у LinkedIn Clone",
        "left": [
            "Work.ua: сильний пошук вакансій, фільтри, резюме та відгуки, але фокус переважно на job-board сценарії.",
            "LinkedIn: сильна глобальна мережа, профіль, контент і вакансії, але продукт перевантажений для локального навчального MVP.",
            "Djinni: зручний IT-рекрутинг і анонімний профіль, але вузький фокус на наймі без широких спільнот і стрічки.",
            "Robota.ua: багато вакансій, резюме і роботодавців, але менше соціальної взаємодії, чатів і професійного контенту.",
        ],
        "right": [
            "У LinkedIn Clone вакансії не ізольовані: вони пов'язані з профілем, чатами, контактами і загальним шляхом кандидата.",
            "LinkedIn Clone бере логіку LinkedIn, але показує її у компактному MVP, де кожен модуль зрозумілий на захисті.",
            "У LinkedIn Clone є не тільки найм, а й контент, groups/pages/events, глобальний пошук і особистий кабінет.",
            "Наш сайт об'єднує job-board, професійну соцмережу і внутрішню комунікацію в одному frontend-сценарії.",
        ],
    },
    {
        "type": "modules",
        "focus": None,
        "title": "Переваги UX нашого сайту",
        "subtitle": "Акцент на зручності інтерфейсу для користувача, а не на backend-інтеграції",
        "footer": "Для кандидата: одна навігація, теми, мови та social demo — все в одному UX.",
        "items": [
            ("Єдиний UX", "одна навігація для всіх сценаріїв"),
            ("Dark/Light", "перемикання теми через CSS variables"),
            ("i18n", "4 мови інтерфейсу в runtime"),
            ("Картки", "зрозуміла структура контенту"),
            ("Анімації", "плавні переходи між сторінками"),
            ("Пошук", "глобальний пошук по сутностях"),
            ("Demo-режим", "social demo без реєстрації по пошті"),
            ("Fallback", "стабільні mock-дані для захисту"),
        ],
    },
    {
        "type": "comparison",
        "focus": None,
        "title": "Що є у конкурентів і що додано у нас",
        "left_title": "У конкурентів",
        "right_title": "У нашому LinkedIn Clone",
        "left": [
            "Вакансії існують окремо від соціальної активності кандидата.",
            "Пошук роботи часто не пов'язаний з чатами та мережею контактів.",
            "Демонстраційні проєкти нерідко залишаються статичними макетами.",
            "Інтерфейс не завжди показує повний шлях користувача від входу до дії.",
        ],
        "right": [
            "Користувач шукає роботу, веде профіль, читає контент і спілкується в одному місці.",
            "Вакансії, контакти та повідомлення об'єднані в логічний кандидатський сценарій.",
            "Frontend має demo-режим і fallback-дані для гарантованої демонстрації на захисті.",
            "На захисті показуємо послідовність через social demo: auth → profile → feed → jobs → chat.",
        ],
    },
    {
        "type": "structure",
        "focus": None,
        "title": "Компонентна структура frontend",
        "subtitle": "Кожен блок — ізольований React-модуль з власною логікою",
        "tree": FRONTEND_TREE_LINES,
        "points": [
            "pages/ — екрани за сценаріями: auth, home, network, jobs, chat, admin",
            "features/ — бізнес-логіка модулів, API-виклики, mappers",
            "app/ — router, providers, layout, ініціалізація UI",
            "shared/ — api, ui, lib: спільний код без прив'язки до одного модуля",
            "contexts — Auth, Profile, Network, Vacancies, Chat (глобальний стан)",
        ],
        "note": "Код розбито за папками так, щоб модулі не залежали один від одного напряму.",
    },
    {
        "type": "architecture",
        "focus": None,
        "title": "Архітектура frontend",
        "steps": [
            ("React/Vite SPA", "UI, routing, компоненти"),
            ("Context Providers", "auth, profile, network, jobs, chat"),
            ("API client", "підготовлений шар запитів"),
            ("Facade.API", "контракт backend-команди"),
            ("Demo + mock", "стабільна демонстрація на захисті"),
        ],
        "note": "Frontend керує станом через Context API та API-клієнт; backend і БД — зона backend-команди, ми показуємо готовий клієнтський шар.",
    },
    {
        "type": "frontend",
        "focus": "home",
        "title": "Маршрутизація frontend",
        "subtitle": "React Router, guarded routes та demo-вхід без live-реєстрації",
        "routes": [
            ("/auth", "публічний маршрут, social demo"),
            ("RequireAuth", "захист приватних сторінок"),
            ("AuthBootstrapGate", "ініціалізація сесії при завантаженні"),
            ("/home … /profile", "приватні маршрути основних модулів"),
            ("PageTransitionOutlet", "анімації переходів між сторінками"),
            ("redirect / → /home", "стартовий сценарій після входу"),
        ],
    },
    {
        "type": "journey",
        "focus": "home",
        "title": "Demo-flow для захисту",
        "steps": [
            "Social demo",
            "Профіль",
            "Стрічка",
            "Спільноти",
            "Вакансії",
            "Чати",
        ],
        "note": "На захисті використовуємо social demo і mock-дані — без live-реєстрації та без залежності від нестабільного API.",
    },
    {
        "type": "admin",
        "focus": None,
        "title": "Адмін-панель платформи",
        "subtitle": "Окремий UI-контур /admin для модерації та керування LinkedIn Clone",
        "access": [
            "/admin/* → RequireAuth + RequireAdmin",
            "JWT role Admin (user.isAdmin)",
            "Без ролі — AdminForbiddenPage (403)",
            "AdminLayout: sidebar з 7 розділами",
            "adminApi.js → /api/admin/*",
        ],
        "sections": [
            ("Dashboard", "stats: users, posts, vacancies, events"),
            ("Users", "list, filters, lock/unlock, roles, create user"),
            ("Content", "moderation постів: search, delete/restore"),
            ("Comments", "moderation коментарів до постів"),
            ("Jobs", "moderation вакансій + recommended queries"),
            ("Events", "moderation подій: delete/restore"),
            ("Roles", "ролі User/Admin у drawer користувача"),
        ],
        "note": "UI адмінки на frontend; demo/mock — порожні дані; повний сценарій — Admin JWT + backend.",
    },
    {
        "type": "roles",
        "focus": None,
        "title": "Розподіл frontend-роботи",
        "left_title": "Ямчук Тимур",
        "right_title": "Андрій Ротарь",
        "left": [
            "Вікно реєстрації та auth UX.",
            "Сторінка вакансій і пошук роботи.",
            "Чати, дизайн і особистий кабінет.",
        ],
        "right": [
            "Головний екран і динамічна контентна стрічка (Feed).",
            "Архітектура Network: контакти, groups, pages, events.",
            "API-клієнт, demo-режим і підготовка до backend-інтеграції.",
        ],
    },
    {
        "type": "feature",
        "focus": "profile",
        "title": "Ямчук Тимур: вікно реєстрації",
        "headline": "Реєстрація є першою точкою входу користувача в платформу.",
        "points": [
            "Побудовано зрозумілий auth UX: користувач бачить форму, стани введення та результат дії.",
            "Опрацьовано валідацію полів і повідомлення про помилки, щоб сценарій не виглядав як статичний макет.",
            "Після успішної реєстрації користувач переходить до подальшого користування сайтом.",
            "Auth-екран стилістично узгоджено з іншими частинами платформи.",
        ],
    },
    {
        "type": "feature",
        "focus": "vacancies",
        "title": "Ямчук Тимур: сторінка вакансій",
        "headline": "Сторінка вакансій реалізує один з головних сценаріїв LinkedIn-like платформи.",
        "points": [
            "Створено список вакансій з картками, основними параметрами і зрозумілою структурою перегляду.",
            "Додано відображення формату роботи, рівня досвіду, локації та інших характеристик вакансії.",
            "Підготовлено сценарії збереження вакансії та відгуку користувача.",
            "Сторінка вписана в загальну навігацію і не виглядає окремим модулем.",
        ],
    },
    {
        "type": "feature",
        "focus": "vacancies",
        "title": "Ямчук Тимур: пошук роботи",
        "headline": "Пошук роботи робить модуль вакансій корисним для щоденного сценарію кандидата.",
        "points": [
            "Реалізовано фільтрацію вакансій за ключовими параметрами.",
            "Опрацьовано сценарій, коли результатів немає або користувач змінює запит.",
            "Пошук допомагає швидко перейти від загального списку до релевантних вакансій.",
            "Цей блок підсилює головну перевагу сайту: вакансії пов'язані з профілем і комунікацією.",
        ],
    },
    {
        "type": "feature",
        "focus": "chat",
        "title": "Ямчук Тимур: чати",
        "headline": "Чати додають платформі внутрішню комунікацію без переходу в сторонні месенджери.",
        "points": [
            "Створено інтерфейс діалогів, повідомлень і станів unread.",
            "Додано сценарії архівації, mute, видалення повідомлень і share posts.",
            "UI-шар для AI-підказок і майбутнього WebRTC — це mock-моделі, не production-функції.",
            "Комунікація підтримує ідею LinkedIn-like продукту: робота, контакти і повідомлення в одному місці.",
        ],
    },
    {
        "type": "feature",
        "focus": "profile",
        "title": "Ямчук Тимур: дизайн і особистий кабінет",
        "headline": "Особистий кабінет і дизайн формують завершене враження від frontend-частини.",
        "points": [
            "Профіль користувача містить особисті та професійні дані, avatar, skills і resume.",
            "Візуальний стиль узгоджено між auth, vacancies, chat і profile.",
            "Кабінет працює як центр професійної ідентичності користувача.",
            "Дизайн побудований на картках, темній/світлій темі, акцентах і зрозумілій навігації.",
        ],
    },
    {
        "type": "feature",
        "focus": "home",
        "title": "Андрій Ротарь: головний екран",
        "headline": "Головний екран задає перше враження і веде користувача до основних розділів платформи.",
        "points": [
            "Створено структуру home-сторінки з фокусом на професійний контент.",
            "Додано швидкий доступ до мережі, вакансій, повідомлень і профілю.",
            "Головний екран працює як стартова точка після авторизації.",
            "Навігація допомагає показати сайт як цілісну платформу.",
        ],
    },
    {
        "type": "feature",
        "focus": "home",
        "title": "Андрій Ротарь: контентна стрічка (Feed)",
        "headline": "Feed — найскладніша frontend-частина: динамічні пости, реакції та оновлення стану.",
        "points": [
            "Реалізовано подачу постів у форматі професійної стрічки з composer і sidebar inbox.",
            "Додано обробку реакцій, коментарів і оновлення UI після дій користувача.",
            "Стрічка синхронізується з providers і mock/API-даними без перезавантаження сторінки.",
            "Це технічно складніший модуль, ніж статичні форми auth або profile.",
        ],
    },
    {
        "type": "feature",
        "focus": "network",
        "title": "Андрій Ротарь: архітектура Network",
        "headline": "Network поєднує контакти, підписки, groups, pages і events в одній вкладці.",
        "points": [
            "Побудовано структуру вкладки network з окремими підрозділами та навігацією.",
            "Реалізовано зв'язки між контактами, профілями та соціальними сутностями.",
            "Groups/pages/events підсилюють відмінність LinkedIn Clone від звичайного job-board.",
            "Модуль показує архітектурну складність frontend, а не лише візуальний дизайн.",
        ],
    },
    {
        "type": "feature",
        "focus": "home",
        "title": "Андрій Ротарь: API-клієнт і demo-режим",
        "headline": "На захисті — demo-режим з fallback; API-клієнт готовий до підключення backend.",
        "points": [
            "Social demo і fallback-дані гарантують стабільну демонстрацію без залежності від сервера.",
            "Створено структуру API-клієнта та proxy для майбутньої backend-інтеграції.",
            "Механізм оновлення сесії є в коді, але live refresh на захисті не показуємо.",
            "Backend-баги — зона backend-команди; frontend-шар готовий до підключення.",
        ],
    },
    {
        "type": "matrix",
        "focus": None,
        "title": "Висновок по frontend-роботі",
        "left_title": "Що реалізовано",
        "right_title": "Практичний результат",
        "left": [
            "Завершений frontend MVP: UI, routing, providers, demo-flow і компонентна архітектура.",
            "Два зони відповідальності: кандидатський UX (Тимур) і контент/мережа/API-клієнт (Андрій).",
            "Demo-режим забезпечує стабільну презентацію без ризику live API-помилок.",
        ],
        "right": [
            "Це не макет, а архітектурно завершений frontend-продукт з реальними сценаріями.",
            "Проєкт показує перевагу LinkedIn-like підходу над простими job-board.",
            "Frontend готовий до подальшого підключення backend після стабілізації API.",
        ],
    },
    {
        "type": "final",
        "focus": None,
        "title": "Підсумок",
        "subtitle": "Frontend MVP з demo-flow — архітектурно завершений продукт для захисту",
        "facts": [
            "LinkedIn Clone — це frontend MVP професійної екосистеми, а не окремий job-board.",
            "На захисті демонструємо demo-flow через social demo та mock-дані.",
            "Головний результат — архітектурно завершений клієнтський шар і зрозумілий UX.",
        ],
    },
]


def apply_report_style():
    report_labels = ["Причина вибору", "Обране рішення", "Реалізація", "Результат"]

    for slide in SLIDES:
        if slide.get("type") == "cover":
            slide["facts"] = [
                "Формат звіту: кожен слайд пояснює причину вибору рішення та результат",
                *slide["facts"],
            ]
            continue

        if not slide["title"].startswith("Звіт:"):
            slide["title"] = f"Звіт: {slide['title']}"

        if slide["type"] == "statement":
            slide["headline"] = f"Обґрунтування рішення: {slide['headline']}"
            slide["points"] = [
                f"{report_labels[i]}: {point}"
                for i, point in enumerate(slide["points"])
            ]

        elif slide["type"] == "filesystem":
            slide["subtitle"] = f"Звітний огляд: {slide['subtitle']}"

        elif slide["type"] == "comparison":
            slide["left_title"] = f"Стан у конкурентів: {slide['left_title']}"
            slide["right_title"] = f"Наше рішення: {slide['right_title']}"
            slide["left"] = [
                f"Аналіз: {point}" for point in slide["left"]
            ]
            slide["right"] = [
                f"Перевага: {point}" for point in slide["right"]
            ]

        elif slide["type"] == "modules":
            slide["subtitle"] = f"Звітний висновок: {slide['subtitle']}"
            slide["items"] = [
                (name, f"Рішення: {desc}") for name, desc in slide["items"]
            ]

        elif slide["type"] == "structure":
            slide["subtitle"] = f"Звітний огляд: {slide['subtitle']}"
            slide["points"] = [f"Модуль: {point}" for point in slide["points"]]
            slide["note"] = f"Архітектура: {slide['note']}"

        elif slide["type"] == "architecture":
            slide["steps"] = [
                (name, f"Чому обрано: {desc}") for name, desc in slide["steps"]
            ]
            slide["note"] = f"Результат рішення: {slide['note']}"

        elif slide["type"] == "frontend":
            slide["routes"] = [
                (route, f"Роль у рішенні: {desc}") for route, desc in slide["routes"]
            ]

        elif slide["type"] == "journey":
            slide["note"] = f"Звітний результат: {slide['note']}"

        elif slide["type"] == "admin":
            slide["subtitle"] = f"Звітний огляд: {slide['subtitle']}"
            slide["note"] = f"Результат: {slide['note']}"

        elif slide["type"] == "roles":
            slide["left"] = [f"Зона відповідальності: {point}" for point in slide["left"]]
            slide["right"] = [f"Зона відповідальності: {point}" for point in slide["right"]]

        elif slide["type"] == "feature":
            slide["headline"] = f"Обґрунтування рішення: {slide['headline']}"
            slide["points"] = [
                f"{report_labels[i]}: {point}"
                for i, point in enumerate(slide["points"])
            ]

        elif slide["type"] == "matrix":
            slide["left_title"] = f"Причина / виконання: {slide['left_title']}"
            slide["right_title"] = f"Результат / висновок: {slide['right_title']}"
            slide["left"] = [f"Звіт: {point}" for point in slide["left"]]
            slide["right"] = [f"Висновок: {point}" for point in slide["right"]]

        elif slide["type"] == "final":
            slide["facts"] = [f"Підсумковий висновок: {fact}" for fact in slide["facts"]]


apply_report_style()
normalize_slide_numbers()

MODULE_ICON_DATA = {
    "Єдиний UX": {"main": "UX", "color_idx": 0},
    "Dark/Light": {"main": "TH", "color_idx": 1},
    "i18n": {"main": "LN", "color_idx": 2},
    "Картки": {"main": "CD", "color_idx": 3},
    "Анімації": {"main": "AN", "color_idx": 4},
    "Пошук": {"main": "SR", "color_idx": 5},
    "Demo-режим": {"main": "DM", "color_idx": 6},
    "Fallback": {"main": "FB", "color_idx": 7},
    "pages/": {"main": "PG", "color_idx": 0},
    "features/": {"main": "FT", "color_idx": 1},
    "app/": {"main": "AP", "color_idx": 2},
    "shared/api": {"main": "API", "color_idx": 3},
    "shared/ui": {"main": "UI", "color_idx": 4},
    "shared/lib": {"main": "LB", "color_idx": 5},
    "contexts": {"main": "CTX", "color_idx": 6},
    "config": {"main": "CFG", "color_idx": 7},
    "Dashboard": {"main": "DB", "color_idx": 0},
    "Users": {"main": "US", "color_idx": 1},
    "Content": {"main": "CT", "color_idx": 2},
    "Comments": {"main": "CM", "color_idx": 3},
    "Jobs": {"main": "JB", "color_idx": 4},
    "Events": {"main": "EV", "color_idx": 5},
    "Roles": {"main": "RL", "color_idx": 6},
}

FEATURE_TASK_META = {
    "вікно реєстрації": {
        "short": "Auth",
        "icon": "A",
        "color_idx": 0,
        "subtitle": "реєстрація, login, ініціалізація сесії",
    },
    "сторінка вакансій": {
        "short": "Jobs",
        "icon": "J",
        "color_idx": 2,
        "subtitle": "список, картки, фільтри, apply",
    },
    "пошук роботи": {
        "short": "Search",
        "icon": "S",
        "color_idx": 4,
        "subtitle": "фільтри, запит, релевантні вакансії",
    },
    "чати": {
        "short": "Chat",
        "icon": "C",
        "color_idx": 5,
        "subtitle": "діалоги, unread, UI mock AI/WebRTC",
    },
    "дизайн і особистий кабінет": {
        "short": "Profile",
        "icon": "P",
        "color_idx": 1,
        "subtitle": "кабінет, avatar, skills, resume",
    },
    "головний екран": {
        "short": "Home",
        "icon": "H",
        "color_idx": 0,
        "subtitle": "стрічка, навігація, швидкий доступ",
    },
    "контентна стрічка": {
        "short": "Feed",
        "icon": "F",
        "color_idx": 3,
        "subtitle": "пости, реакції, коментарі, providers",
    },
    "архітектура network": {
        "short": "Network",
        "icon": "N",
        "color_idx": 4,
        "subtitle": "контакти, groups, pages, events",
    },
    "злиття backend з frontend": {
        "short": "API",
        "icon": "API",
        "color_idx": 7,
        "subtitle": "API-клієнт, demo-режим, fallback",
    },
    "api-клієnt і demo-режим": {
        "short": "API",
        "icon": "API",
        "color_idx": 7,
        "subtitle": "API-клієnt, demo-режим, fallback",
    },
}

MODULE_COLORS = [BLUE, PURPLE, CYAN, GREEN, ORANGE, BLUE, PURPLE, CYAN]

NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"


def add_shape(slide, shape_type, x, y, w, h, color, line_color=None):
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.fill.solid()
        shape.line.color.rgb = line_color
    return shape


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = Inches(0.04)
    box.text_frame.margin_right = Inches(0.04)
    box.text_frame.margin_top = Inches(0.03)
    box.text_frame.margin_bottom = Inches(0.03)
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_centered_text(slide, x, y, w, h, text, size=18, bold=False, color=TEXT):
    box = add_text(slide, x, y, w, h, text, size=size, bold=bold, color=color, align=PP_ALIGN.CENTER)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return box


def add_bullet_list(slide, items, x, y, w, line_h=0.56, size=13, color=TEXT):
    for i, txt in enumerate(items):
        yy = y + i * line_h
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, x, yy + 0.1, 0.12, 0.12, PURPLE_SOFT)
        add_text(slide, x + 0.25, yy, w - 0.3, line_h - 0.04, txt, size=size, color=color)


def add_title(slide, title, subtitle=None):
    add_text(slide, 0.7, 0.86, 12.0, 0.55, title, size=24, bold=True, color=TEXT)
    if subtitle:
        add_text(slide, 0.72, 1.38, 11.4, 0.32, subtitle, size=12, color=TEXT_MUTED)


def add_card_with_shadow(slide, x, y, w, h, radius_color=CARD):
    shadow = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + 0.04, y + 0.05, w, h, RGBColor(7, 9, 13))
    card = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h, radius_color, line_color=BORDER)
    return shadow, card


def draw_icon_badge(slide, x, y, size, icon, bg_color, font_size=14, text_color=WHITE):
    if len(icon) > 1:
        w = size * 1.25
        shape = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, size, bg_color, line_color=bg_color)
        font_size = min(font_size, 9)
    else:
        w = size
        shape = add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, x, y, size, size, bg_color, line_color=bg_color)
    tf = shape.text_frame
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = icon
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = text_color
    return shape


def track_shape(shape, anim_ids):
    if anim_ids is not None:
        anim_ids.append(shape.shape_id)
    return shape


def add_click_appear_sequence(slide, shape_ids, duration_ms=450):
    if not shape_ids:
        return

    effect_nodes = []
    node_id = 3
    for spid in shape_ids:
        effect_nodes.append(
            f"""
            <p:par>
              <p:cTn id="{node_id}" presetID="10" presetClass="entr" presetSubtype="0"
                     fill="hold" nodeType="clickEffect">
                <p:stCondLst>
                  <p:cond delay="indefinite"/>
                </p:stCondLst>
                <p:childTnLst>
                  <p:animEffect transition="in" filter="fade">
                    <p:cBhvr>
                      <p:cTn id="{node_id + 1}" dur="{duration_ms}"/>
                      <p:tgtEl>
                        <p:spTgt spid="{spid}"/>
                      </p:tgtEl>
                    </p:cBhvr>
                  </p:animEffect>
                </p:childTnLst>
              </p:cTn>
            </p:par>"""
        )
        node_id += 2

    timing_xml = f"""
    <p:timing xmlns:p="{NS_P}">
      <p:tnLst>
        <p:par>
          <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
            <p:childTnLst>
              <p:seq concurrent="1" nextAc="seek">
                <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
                  <p:childTnLst>
                    {''.join(effect_nodes)}
                  </p:childTnLst>
                </p:cTn>
              </p:seq>
            </p:childTnLst>
          </p:cTn>
        </p:par>
      </p:tnLst>
    </p:timing>
    """
    slide._element.append(parse_xml(timing_xml))


def get_module_icons(name):
    return MODULE_ICON_DATA.get(name, {"main": "◆", "color_idx": 0})


def get_feature_task_label(data):
    title = data.get("title", "")
    title = re.sub(r"^\d+\.\s*", "", title)
    title = re.sub(r"^Звіт:\s*", "", title)
    parts = [part.strip() for part in title.split(":") if part.strip()]
    if len(parts) >= 2:
        return parts[-1]
    return title


def get_feature_task_meta(data):
    label = get_feature_task_label(data).lower()
    for key, meta in FEATURE_TASK_META.items():
        if key in label:
            return meta
    short = label.split()[0][:10].title() if label else "Task"
    return {
        "short": short,
        "icon": short[0].upper(),
        "color_idx": 0,
        "subtitle": label,
    }


def draw_green_check(slide, x, y, size=16):
    return add_centered_text(slide, x, y, 0.3, 0.3, "✓", size=size, bold=True, color=GREEN)


def add_check_list(slide, items, x, y, w, line_h=0.82, size=9.5, color=TEXT):
    for i, txt in enumerate(items):
        yy = y + i * line_h
        draw_green_check(slide, x, yy + 0.04, size=14)
        add_text(slide, x + 0.28, yy, w - 0.32, line_h - 0.04, txt, size=size, color=color)


def draw_task_card(slide, x, y, w, h, meta, anim_ids=None):
    color = MODULE_COLORS[meta.get("color_idx", 0) % len(MODULE_COLORS)]
    icon_text = meta.get("icon", "T")
    header_y = y + 0.14
    icon_size = 0.4
    add_card_with_shadow(slide, x, y, w, h, CARD)
    draw_icon_badge(slide, x + 0.16, header_y, icon_size, icon_text, color, font_size=12)
    draw_green_check(slide, x + w - 0.4, header_y + 0.05, size=17)
    title_w = w - 1.05
    add_centered_text(slide, x + 0.62, header_y + 0.04, title_w, 0.32, meta["short"], size=13, bold=True, color=TEXT)
    add_text(slide, x + 0.16, y + 0.58, w - 0.32, 0.46, meta.get("subtitle", ""), size=9.5, color=TEXT_MUTED)
    return meta


def draw_number_circle(slide, x, y, size, number, color):
    shape = add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, x, y, size, size, color)
    tf = shape.text_frame
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = str(number)
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(17 if size < 1.0 else 20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    return shape


def draw_nav_buttons(slide, focus):
    start_x = 2.55
    y = 0.14
    w = 1.62
    gap = 0.12
    for i, (key, label) in enumerate(NAV_ITEMS):
        x = start_x + i * (w + gap)
        active = key == focus
        shape = add_shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            x,
            y,
            w,
            0.28,
            PURPLE_SOFT if active else CARD_SOFT,
            line_color=BORDER,
        )
        p = shape.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(10)
        p.font.bold = active
        p.font.color.rgb = RGBColor(20, 24, 36) if active else TEXT_MUTED
        p.alignment = PP_ALIGN.CENTER


def draw_base(slide, idx, total, focus, show_nav=True, show_header=True):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 7.5, BG_DARK)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, 10.55, 0.35, 2.5, 2.5, RGBColor(18, 32, 55))
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, 0.35, 5.45, 1.55, 1.55, RGBColor(22, 24, 47))
    if show_header:
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 0.72, BG_DARK_2)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0.7, 13.333, 0.02, BORDER)
        logo = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.62, 0.13, 0.46, 0.42, BLUE, line_color=BLUE)
        logo.text_frame.text = "in"
        p = logo.text_frame.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        add_text(slide, 1.18, 0.19, 1.55, 0.25, "LinkedIn Clone", size=11, bold=True, color=TEXT)
        if show_nav:
            draw_nav_buttons(slide, focus)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 7.32, 13.333, 0.18, BG_DARK_2)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 7.32, 13.333 * idx / total, 0.18, BLUE)
    add_text(slide, 0.75, 7.31, 2.3, 0.12, "LinkedIn Diplom", size=8, color=TEXT_MUTED)
    add_text(slide, 12.3, 7.31, 0.9, 0.12, f"{idx}/{total}", size=8, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)


def draw_cover(slide, data, idx, total):
    draw_base(slide, idx, total, data.get("focus"), show_nav=False, show_header=False)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, 10.8, 0.55, 2.0, 2.0, BLUE_SOFT)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, 10.2, 2.35, 2.4, 2.4, RGBColor(17, 37, 63))
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.68, 0.95, 1.18, 0.9, BLUE, line_color=BLUE).text_frame.text = "in"
    p = slide.shapes[-1].text_frame.paragraphs[0]
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    add_text(slide, 2.15, 1.02, 9.45, 0.9, data["title"], size=39, bold=True, color=WHITE)
    add_text(slide, 2.18, 1.9, 9.1, 0.68, data["subtitle"], size=19, color=RGBColor(203, 213, 225))
    add_card_with_shadow(slide, 2.18, 2.95, 9.9, 3.35, CARD)
    for i, fact in enumerate(data["facts"]):
        add_text(slide, 2.55, 3.22 + i * 0.56, 9.05, 0.48, fact, size=14, color=TEXT)


def draw_cards(slide, data, idx, total):
    draw_base(slide, idx, total, data.get("focus"))
    add_title(slide, data["title"], "Ключові аргументи та результати")
    cards = [(0.8, 1.98), (6.95, 1.98), (0.8, 4.33), (6.95, 4.33)]
    for i, txt in enumerate(data["points"]):
        x, y = cards[i]
        add_card_with_shadow(slide, x, y, 5.55, 2.1, CARD)
        badge = add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, x + 0.2, y + 0.2, 0.45, 0.45, BLUE)
        badge.text_frame.text = str(i + 1)
        bp = badge.text_frame.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        bp.font.bold = True
        bp.font.size = Pt(12)
        bp.font.color.rgb = WHITE
        add_text(slide, x + 0.78, y + 0.22, 4.55, 1.72, txt, size=12.5, color=TEXT)


def draw_matrix(slide, data, idx, total):
    draw_base(slide, idx, total, data.get("focus"))
    add_title(slide, data["title"])
    add_card_with_shadow(slide, 0.8, 1.95, 5.95, 5.05, CARD)
    add_card_with_shadow(slide, 6.55, 1.95, 5.95, 5.05, CARD)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0.8, 1.95, 5.95, 0.05, PURPLE_SOFT)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 6.55, 1.95, 5.95, 0.05, PURPLE_SOFT)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 1.05, 2.18, 2.45, 0.38, BLUE_SOFT)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 6.8, 2.18, 2.95, 0.38, BLUE_SOFT)
    add_text(slide, 1.2, 2.25, 2.1, 0.2, data.get("left_title", "Факти"), size=11, bold=True, color=WHITE)
    add_text(slide, 6.95, 2.25, 2.9, 0.2, data.get("right_title", "Практичні висновки"), size=11, bold=True, color=WHITE)
    left_items = data["left"]
    right_items = data["right"]
    left_gap = 4.15 / max(1, len(left_items))
    right_gap = 4.15 / max(1, len(right_items))
    for i, txt in enumerate(left_items):
        y = 2.6 + i * left_gap
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 1.02, y, 5.5, left_gap - 0.1, CARD_SOFT, line_color=BORDER)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 1.08, y + 0.14, 0.07, left_gap - 0.38, PURPLE_SOFT)
        add_text(slide, 1.28, y + 0.08, 5.05, left_gap - 0.16, txt, size=11.5, color=TEXT)
    for i, txt in enumerate(right_items):
        y = 2.6 + i * right_gap
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 6.77, y, 5.5, right_gap - 0.1, CARD_SOFT, line_color=BORDER)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 6.83, y + 0.14, 0.07, right_gap - 0.38, PURPLE_SOFT)
        add_text(slide, 7.03, y + 0.08, 5.05, right_gap - 0.16, txt, size=11.5, color=TEXT)


def draw_process(slide, data, idx, total):
    draw_base(slide, idx, total, data.get("focus"))
    add_title(slide, data["title"], "Поетапна логіка реалізації")
    steps = data["steps"]
    count = max(1, len(steps))
    grid_start = 0.85
    grid_width = 11.65
    gap = 0.28
    step_w = (grid_width - gap * (count - 1)) / count
    for i, step in enumerate(steps):
        x = grid_start + i * (step_w + gap)
        color = CARD if i % 2 == 0 else CARD_SOFT
        add_card_with_shadow(slide, x, 2.2, step_w, 2.65, color)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, x + (step_w / 2) - 0.27, 2.45, 0.55, 0.55, BLUE).text_frame.text = str(i + 1)
        p = slide.shapes[-1].text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        add_text(slide, x + 0.12, 3.16, step_w - 0.24, 1.4, step, size=12, align=PP_ALIGN.CENTER, color=TEXT)
        if i < count - 1:
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.CHEVRON, x + step_w + 0.03, 3.25, 0.18, 0.5, BLUE_SOFT)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.9, 5.45, 11.95, 1.2, CARD)
    add_text(slide, 1.2, 5.77, 11.3, 0.65, "Кожний етап має пряме технічне підтвердження у коді, маршрутах API та робочих екранах frontend.", size=13, color=TEXT_MUTED)


def draw_two_col(slide, data, idx, total):
    draw_base(slide, idx, total, data.get("focus"))
    add_title(slide, data["title"])
    add_card_with_shadow(slide, 0.8, 1.85, 5.85, 5.25, CARD)
    add_card_with_shadow(slide, 6.65, 1.85, 5.85, 5.25, CARD)
    add_text(slide, 1.1, 2.18, 5.1, 0.45, data["col1_title"], size=16, bold=True, color=TEXT)
    add_text(slide, 6.95, 2.18, 5.1, 0.45, data["col2_title"], size=16, bold=True, color=TEXT)
    for i, val in enumerate(data["col1"]):
        add_text(slide, 1.1, 2.7 + i * 0.78, 5.1, 0.65, f"• {val}", size=13, color=TEXT)
    for i, val in enumerate(data["col2"]):
        add_text(slide, 6.95, 2.7 + i * 0.78, 5.1, 0.65, f"• {val}", size=13, color=TEXT)


def draw_report(slide, data, idx, total):
    draw_base(slide, idx, total, data.get("focus"))
    add_title(slide, data["title"], "Звітні твердження з підтвердженням")
    for i, row in enumerate(data["rows"]):
        y = 1.95 + i * 1.33
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.9, y, 3.4, 1.02, BLUE_SOFT)
        add_card_with_shadow(slide, 4.45, y, 7.95, 1.02, CARD)
        add_text(slide, 1.2, y + 0.24, 2.9, 0.55, row[0], size=13, bold=True, color=WHITE)
        add_text(slide, 4.75, y + 0.2, 7.4, 0.62, row[1], size=13, color=TEXT)


def draw_statement(slide, data, idx, total):
    draw_base(slide, idx, total, data.get("focus"))
    add_title(slide, data["title"])
    add_text(slide, 0.9, 1.68, 7.0, 1.75, data["headline"], size=21, bold=True, color=WHITE)
    add_card_with_shadow(slide, 8.65, 1.55, 3.15, 3.0, CARD_SOFT)
    add_centered_text(slide, 9.0, 1.95, 2.45, 1.15, data["metric"], size=70, bold=True, color=PURPLE_SOFT)
    add_centered_text(slide, 9.0, 3.15, 2.45, 0.7, data["metric_label"], size=13, color=TEXT)
    add_card_with_shadow(slide, 0.95, 3.78, 7.35, 2.12, CARD)
    add_bullet_list(slide, data["points"], 1.22, 4.02, 6.5, line_h=0.56, size=13)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 8.95, 4.95, 2.55, 0.55, BLUE, line_color=BLUE)
    add_centered_text(slide, 9.04, 5.08, 2.35, 0.25, "MVP для захисту", size=12, bold=True, color=WHITE)


def draw_modules(slide, data, idx, total, anim_ids=None):
    draw_base(slide, idx, total, data.get("focus"))
    add_title(slide, data["title"], data.get("subtitle"))
    cols = 4
    card_w = 2.86
    card_h = 1.28
    start_x = 0.82
    start_y = 1.95
    gap_x = 0.32
    gap_y = 0.32
    colors = MODULE_COLORS
    for i, (name, desc) in enumerate(data["items"]):
        row = i // cols
        col = i % cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        icon_meta = get_module_icons(name)
        color = colors[icon_meta.get("color_idx", i) % len(colors)]
        add_card_with_shadow(slide, x, y, card_w, card_h, CARD)
        main_icon = draw_icon_badge(
            slide,
            x + 0.16,
            y + 0.18,
            0.42,
            icon_meta["main"],
            color,
            font_size=15,
        )
        track_shape(main_icon, anim_ids)
        add_text(slide, x + 0.68, y + 0.2, card_w - 0.88, 0.3, name, size=13, bold=True, color=TEXT)
        add_text(slide, x + 0.2, y + 0.72, card_w - 0.4, 0.42, desc, size=10.5, color=TEXT_MUTED)
    add_card_with_shadow(slide, 1.15, 5.5, 11.0, 0.82, CARD_SOFT)
    footer = data.get(
        "footer",
        "Єдина платформа: користувач не перемикається між окремими сервісами для роботи, контенту та комунікації.",
    )
    add_centered_text(slide, 1.35, 5.66, 10.6, 0.42, footer, size=13, color=TEXT)


def draw_structure(slide, data, idx, total):
    draw_base(slide, idx, total, data.get("focus"))
    add_title(slide, data["title"], data.get("subtitle"))
    add_card_with_shadow(slide, 0.78, 1.78, 5.55, 4.55, CARD)
    add_text(slide, 1.02, 2.0, 4.8, 0.28, "Дерево папок frontend", size=13, bold=True, color=CYAN)
    for i, line in enumerate(data.get("tree", FRONTEND_TREE_LINES)):
        color = TEXT if line.endswith("/") or line.startswith("frontend") else TEXT_MUTED
        add_text(slide, 1.05, 2.38 + i * 0.34, 5.1, 0.3, line, size=9.5, color=color)
    add_card_with_shadow(slide, 6.62, 1.78, 6.0, 4.55, CARD_SOFT)
    add_text(slide, 6.88, 2.0, 5.4, 0.28, "Ізольовані React-модулі", size=13, bold=True, color=PURPLE_SOFT)
    add_check_list(slide, data["points"], 6.88, 2.42, 5.5, line_h=0.72, size=10, color=TEXT)
    add_card_with_shadow(slide, 0.78, 6.05, 11.85, 0.62, CARD_SOFT)
    add_text(slide, 1.02, 6.18, 11.4, 0.38, data["note"], size=10.5, color=TEXT)


def draw_filesystem_tree_panel(slide, x, y, w, h, label, lines, label_color):
    add_card_with_shadow(slide, x, y, w, h, CARD)
    add_text(slide, x + 0.18, y + 0.12, w - 0.36, 0.24, label, size=12, bold=True, color=label_color)
    inner_y = y + 0.42
    for i, line in enumerate(lines):
        is_root = line.endswith("/") or "frontend/" in line or "backend/" in line
        add_text(
            slide,
            x + 0.16,
            inner_y + i * 0.31,
            w - 0.32,
            0.28,
            line,
            size=9.2 if not is_root else 10,
            bold=is_root,
            color=TEXT if is_root else TEXT_MUTED,
        )


def draw_filesystem(slide, data, idx, total):
    draw_base(slide, idx, total, data.get("focus"))
    add_title(slide, data["title"], data.get("subtitle"))
    draw_filesystem_tree_panel(slide, 0.78, 1.72, 7.35, 4.95, "frontend/ — наш фокус", FRONTEND_TREE_LINES, CYAN)
    draw_filesystem_tree_panel(slide, 8.35, 1.72, 4.58, 4.95, "backend/ — контекст", BACKEND_TREE_LINES, ORANGE)
    add_card_with_shadow(slide, 0.78, 6.05, 12.15, 0.62, CARD_SOFT)
    add_text(
        slide,
        1.02,
        6.18,
        11.65,
        0.38,
        "Структура з реального репозиторію: React SPA розбито на pages, features, app та shared.",
        size=10.5,
        color=TEXT,
    )


def draw_comparison(slide, data, idx, total):
    draw_base(slide, idx, total, data.get("focus"))
    add_title(slide, data["title"])
    add_card_with_shadow(slide, 0.78, 1.78, 5.82, 4.95, CARD)
    add_card_with_shadow(slide, 6.72, 1.78, 5.82, 4.95, CARD_SOFT)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0.78, 1.78, 5.82, 0.07, ORANGE)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 6.72, 1.78, 5.82, 0.07, GREEN)
    add_text(slide, 1.15, 2.02, 4.45, 0.82, data["left_title"], size=14, bold=True, color=ORANGE)
    add_text(slide, 7.08, 2.02, 4.45, 0.82, data["right_title"], size=14, bold=True, color=GREEN)
    add_bullet_list(slide, data["left"], 1.2, 2.95, 4.35, line_h=0.92, size=11.5)
    add_bullet_list(slide, data["right"], 7.13, 2.95, 4.35, line_h=0.92, size=11.5)


def add_fitted_picture(slide, image_path, x, y, max_w, max_h):
    pic = slide.shapes.add_picture(str(image_path), Inches(x), Inches(y))
    max_w_emu = Inches(max_w)
    max_h_emu = Inches(max_h)
    scale = min(max_w_emu / pic.width, max_h_emu / pic.height)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = int(Inches(x) + (max_w_emu - pic.width) / 2)
    pic.top = int(Inches(y) + (max_h_emu - pic.height) / 2)
    return pic


def draw_stack(slide, data, idx, total):
    draw_base(slide, idx, total, data.get("focus"))
    add_title(slide, data["title"], "Стек відповідає реальному коду проєкту та документації")
    for i, (group, items) in enumerate(data["groups"]):
        x = 0.82 + (i % 2) * 6.0
        y = 1.9 + (i // 2) * 2.38
        add_card_with_shadow(slide, x, y, 5.65, 2.05, CARD)
        add_text(slide, x + 0.28, y + 0.22, 5.0, 0.32, group, size=17, bold=True, color=PURPLE_SOFT)
        for j, item in enumerate(items):
            chip_x = x + 0.28 + (j % 3) * 1.72
            chip_y = y + 0.75 + (j // 3) * 0.52
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, chip_x, chip_y, 1.55, 0.33, CARD_LIGHT, line_color=BORDER)
            add_centered_text(slide, chip_x + 0.04, chip_y + 0.06, 1.47, 0.16, item, size=8.5, color=TEXT)


def draw_architecture(slide, data, idx, total):
    draw_base(slide, idx, total, data.get("focus"))
    add_title(slide, data["title"], "Frontend: Context API, API-клієнт, demo-режим")
    count = len(data["steps"])
    start_x = 0.72
    y = 2.08
    w = 2.25
    gap = 0.28
    colors = [BLUE, PURPLE, CYAN, GREEN, ORANGE]
    for i, (name, desc) in enumerate(data["steps"]):
        x = start_x + i * (w + gap)
        add_card_with_shadow(slide, x, y, w, 2.05, CARD)
        add_centered_text(slide, x + 0.22, y + 0.3, w - 0.44, 0.5, name, size=14, bold=True, color=colors[i])
        add_centered_text(slide, x + 0.2, y + 0.95, w - 0.4, 0.55, desc, size=11, color=TEXT)
        if i < count - 1:
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.CHEVRON, x + w + 0.05, y + 0.75, 0.18, 0.55, BLUE_SOFT)
    add_card_with_shadow(slide, 1.1, 5.08, 11.1, 1.05, CARD_SOFT)
    add_text(slide, 1.45, 5.36, 10.5, 0.45, data["note"], size=13, color=TEXT)


def draw_database(slide, data, idx, total):
    draw_base(slide, idx, total, data.get("focus"))
    add_title(slide, data["title"], "Логічна ізоляція даних за доменами")
    add_card_with_shadow(slide, 0.92, 1.86, 4.1, 4.45, CARD)
    add_centered_text(slide, 1.35, 2.18, 3.25, 0.42, "PostgreSQL 16", size=20, bold=True, color=CYAN)
    cols = 3
    for i, schema in enumerate(data["schemas"]):
        x = 1.2 + (i % cols) * 1.18
        y = 2.9 + (i // cols) * 0.72
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, 1.0, 0.42, CARD_LIGHT, line_color=BORDER)
        add_centered_text(slide, x + 0.03, y + 0.1, 0.94, 0.16, schema, size=7.7, color=TEXT)
    add_card_with_shadow(slide, 5.55, 1.86, 6.85, 4.45, CARD_SOFT)
    add_bullet_list(slide, data["points"], 5.9, 2.28, 6.0, line_h=0.78, size=13)


def draw_security(slide, data, idx, total):
    draw_base(slide, idx, total, data.get("focus"))
    add_title(slide, data["title"])
    for i, (name, desc) in enumerate(data["points"]):
        x = 0.9 + (i % 2) * 5.9
        y = 1.88 + (i // 2) * 2.08
        add_card_with_shadow(slide, x, y, 5.45, 1.72, CARD)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, x + 0.25, y + 0.42, 0.62, 0.62, BLUE if i % 2 == 0 else PURPLE)
        add_centered_text(slide, x + 0.25, y + 0.53, 0.62, 0.26, "✓", size=17, bold=True, color=WHITE)
        add_text(slide, x + 1.1, y + 0.28, 4.0, 0.35, name, size=16, bold=True, color=TEXT)
        add_text(slide, x + 1.1, y + 0.72, 4.0, 0.55, desc, size=12, color=TEXT_MUTED)


def draw_frontend(slide, data, idx, total):
    draw_base(slide, idx, total, data.get("focus"))
    add_title(slide, data["title"], data.get("subtitle", "React Router, guarded routes, demo-вхід"))
    add_card_with_shadow(slide, 0.86, 1.78, 11.65, 4.95, CARD)
    for i, (route, desc) in enumerate(data["routes"]):
        x = 1.22 + (i % 3) * 3.62
        y = 2.16 + (i // 3) * 2.0
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, 3.1, 1.48, CARD_SOFT, line_color=BORDER)
        add_text(slide, x + 0.24, y + 0.25, 2.65, 0.28, route, size=18, bold=True, color=PURPLE_SOFT)
        add_text(slide, x + 0.24, y + 0.72, 2.65, 0.44, desc, size=11.5, color=TEXT)


def draw_admin(slide, data, idx, total):
    draw_base(slide, idx, total, data.get("focus"))
    add_title(slide, data["title"], data.get("subtitle"))
    add_card_with_shadow(slide, 0.78, 1.78, 3.48, 4.12, CARD)
    add_text(slide, 1.0, 1.98, 3.0, 0.28, "Доступ і архітектура", size=13, bold=True, color=PURPLE_SOFT)
    add_check_list(slide, data["access"], 1.0, 2.38, 3.05, line_h=0.74, size=9, color=TEXT)
    add_card_with_shadow(slide, 4.42, 1.78, 8.52, 4.12, CARD_SOFT)
    add_text(slide, 4.62, 1.98, 4.5, 0.28, "Розділи admin-панелі", size=13, bold=True, color=CYAN)
    colors = MODULE_COLORS
    row_h = 0.52
    start_y = 2.38
    for i, (name, desc) in enumerate(data["sections"]):
        y = start_y + i * row_h
        icon_meta = get_module_icons(name)
        color = colors[icon_meta.get("color_idx", i) % len(colors)]
        if i > 0:
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 4.58, y - 0.02, 8.18, 0.01, BORDER)
        draw_icon_badge(slide, 4.58, y + 0.06, 0.3, icon_meta["main"], color, font_size=9)
        add_text(slide, 4.98, y + 0.08, 1.35, 0.22, name, size=11.5, bold=True, color=TEXT)
        add_text(slide, 6.35, y + 0.08, 6.2, 0.36, desc, size=9, color=TEXT_MUTED)
        draw_green_check(slide, 12.62, y + 0.08, size=14)
    add_card_with_shadow(slide, 0.78, 6.02, 12.15, 0.58, CARD_SOFT)
    add_text(slide, 1.0, 6.14, 11.7, 0.34, data["note"], size=10, color=TEXT)


def draw_journey(slide, data, idx, total):
    draw_base(slide, idx, total, data.get("focus"))
    add_title(slide, data["title"], "Social demo, mock-дані, без live API")
    steps = data["steps"]
    count = len(steps)
    circle_d = 0.9
    start_x = 0.9
    lane_w = 11.55
    step_pitch = (lane_w - circle_d) / max(1, count - 1)
    y = 2.55
    label_w = max(1.35, step_pitch + 0.15)
    for i, step in enumerate(steps):
        x = start_x + i * step_pitch
        draw_number_circle(slide, x, y, circle_d, i + 1, BLUE if i < 3 else PURPLE)
        add_centered_text(slide, x - (label_w - circle_d) / 2, y + circle_d + 0.16, label_w, 0.58, step, size=11, color=TEXT)
        if i < count - 1:
            chevron_x = x + circle_d + 0.06
            chevron_w = max(0.18, step_pitch - circle_d - 0.12)
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.CHEVRON, chevron_x, y + 0.28, chevron_w, 0.34, BLUE_SOFT)
    add_card_with_shadow(slide, 1.25, 5.25, 10.75, 0.9, CARD_SOFT)
    add_centered_text(slide, 1.55, 5.42, 10.15, 0.42, data["note"], size=13, color=TEXT)


def draw_feature(slide, data, idx, total, anim_ids=None):
    draw_base(slide, idx, total, data.get("focus"))
    add_title(slide, data["title"])
    task_meta = get_feature_task_meta(data)
    card_w = 3.15
    card_h = 1.1
    card_x = 13.333 - 0.68 - card_w
    card_y = 1.5
    draw_task_card(slide, card_x, card_y, card_w, card_h, task_meta, anim_ids)
    add_text(slide, 0.92, 1.75, card_x - 1.0, 0.92, data["headline"], size=22, bold=True, color=WHITE)
    add_card_with_shadow(slide, 0.95, 2.95, 11.45, 2.95, CARD)
    add_bullet_list(slide, data["points"], 1.26, 3.28, 10.8, line_h=0.58, size=12.8)


def draw_integration(slide, data, idx, total):
    draw_base(slide, idx, total, data.get("focus"))
    add_title(slide, data["title"], "Єдиний клієнт API, токени в localStorage, refresh при 401")
    for i, (domain, routes) in enumerate(data["rows"]):
        y = 1.82 + i * 0.72
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.9, y, 2.35, 0.48, BLUE_SOFT, line_color=BORDER)
        add_card_with_shadow(slide, 3.45, y - 0.02, 8.35, 0.52, CARD)
        add_centered_text(slide, 1.0, y + 0.1, 2.15, 0.16, domain, size=10.5, bold=True, color=WHITE)
        add_text(slide, 3.7, y + 0.08, 7.8, 0.2, routes, size=10.3, color=TEXT)


def draw_quality(slide, data, idx, total):
    draw_base(slide, idx, total, data.get("focus"))
    add_title(slide, data["title"])
    for i, (name, desc) in enumerate(data["items"]):
        x = 0.88 + (i % 3) * 4.0
        y = 1.92 + (i // 3) * 2.0
        add_card_with_shadow(slide, x, y, 3.55, 1.52, CARD)
        add_text(slide, x + 0.24, y + 0.22, 3.0, 0.3, name, size=15, bold=True, color=PURPLE_SOFT)
        add_text(slide, x + 0.24, y + 0.72, 3.05, 0.42, desc, size=11.5, color=TEXT)


def draw_infrastructure(slide, data, idx, total):
    draw_base(slide, idx, total, data.get("focus"))
    add_title(slide, data["title"], "Проєкт відтворюється локально та через Docker Compose")
    for i, (name, desc) in enumerate(data["points"]):
        y = 1.95 + i * 1.02
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 1.05, y, 2.65, 0.68, CARD_LIGHT, line_color=BORDER)
        add_card_with_shadow(slide, 4.0, y - 0.02, 7.65, 0.72, CARD)
        add_centered_text(slide, 1.18, y + 0.18, 2.38, 0.18, name, size=12, bold=True, color=CYAN)
        add_text(slide, 4.3, y + 0.18, 7.05, 0.22, desc, size=12, color=TEXT)


def draw_roadmap(slide, data, idx, total):
    draw_base(slide, idx, total, data.get("focus"))
    add_title(slide, data["title"], "Що чесно винесено за межі першої версії")
    for i, (name, desc) in enumerate(data["steps"]):
        x = 0.9 + i * 2.45
        add_card_with_shadow(slide, x, 2.05, 2.0, 3.55, CARD)
        add_centered_text(slide, x + 0.32, 2.38, 1.35, 0.52, str(i + 1), size=28, bold=True, color=PURPLE_SOFT)
        add_centered_text(slide, x + 0.16, 3.18, 1.68, 0.34, name, size=13, bold=True, color=TEXT)
        add_centered_text(slide, x + 0.18, 3.82, 1.64, 0.82, desc, size=10.5, color=TEXT_MUTED)


def draw_roles(slide, data, idx, total):
    draw_matrix(slide, data, idx, total)


def draw_final(slide, data, idx, total):
    draw_base(slide, idx, total, data.get("focus"))
    add_text(slide, 1.0, 1.35, 11.0, 0.8, data["title"], size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    subtitle = data.get(
        "subtitle",
        "Frontend MVP з demo-flow — архітектурно завершений продукт для захисту",
    )
    add_text(slide, 1.35, 2.15, 10.65, 0.45, subtitle, size=17, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    for i, fact in enumerate(data["facts"]):
        add_card_with_shadow(slide, 1.55, 3.05 + i * 0.9, 10.2, 0.68, CARD_SOFT)
        add_text(slide, 1.9, 3.22 + i * 0.9, 9.45, 0.24, fact, size=13.5, color=TEXT)


def build_presentation():
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    total = len(SLIDES)

    draw_map = {
        "cover": draw_cover,
        "cards": draw_cards,
        "matrix": draw_matrix,
        "process": draw_process,
        "two_col": draw_two_col,
        "report": draw_report,
        "statement": draw_statement,
        "modules": draw_modules,
        "structure": draw_structure,
        "comparison": draw_comparison,
        "filesystem": draw_filesystem,
        "stack": draw_stack,
        "architecture": draw_architecture,
        "database": draw_database,
        "security": draw_security,
        "frontend": draw_frontend,
        "journey": draw_journey,
        "admin": draw_admin,
        "feature": draw_feature,
        "integration": draw_integration,
        "quality": draw_quality,
        "infrastructure": draw_infrastructure,
        "roadmap": draw_roadmap,
        "roles": draw_roles,
        "final": draw_final,
    }

    for idx, data in enumerate(SLIDES, start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        if data["type"] == "modules":
            anim_ids = []
            draw_map[data["type"]](slide, data, idx, total, anim_ids)
            add_click_appear_sequence(slide, anim_ids)
        elif data["type"] == "feature":
            draw_map[data["type"]](slide, data, idx, total, None)
        else:
            draw_map[data["type"]](slide, data, idx, total)

    presentation.save(OUTPUT_DOCS_FILE)
    presentation.save(OUTPUT_ROOT_FILE)


if __name__ == "__main__":
    build_presentation()
    print(f"Presentation saved: {OUTPUT_DOCS_FILE}")
    print(f"Presentation saved: {OUTPUT_ROOT_FILE}")
